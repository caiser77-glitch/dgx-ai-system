from pathlib import Path
from pptx import Presentation

output_dir = Path("outputs")
output_dir.mkdir(exist_ok=True)

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "DGX 자동 생성 PPT"
slide.placeholders[1].text = "Python으로 생성한 테스트 프레젠테이션입니다."

out = output_dir / "test_presentation.pptx"
prs.save(out)

print(f"PPT 생성 완료: {out}")
