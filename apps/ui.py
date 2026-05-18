import shutil
import subprocess
import requests
from pathlib import Path

import gradio as gr
from pptx import Presentation

BASE = Path("/home/caiser77/dgx_workspace")

OLLAMA_URL = "http://localhost:11434/api/generate"
OPEN_WEBUI_URL = "http://100.98.149.128:3000/"

MODEL_FAST = "qwen2.5:14b"
MODEL_BIG = "qwen2.5:72b"
MODEL_IMAGE = "llava:latest"
MODEL_LIGHT = "mistral:latest"


def call_ollama(model, prompt, temperature=0.2, timeout=600):
    r = requests.post(
        OLLAMA_URL,
        json={
            "model": model,
            "prompt": prompt,
            "stream": False,
            "options": {"temperature": temperature, "top_p": 0.8},
        },
        timeout=timeout,
    )
    r.raise_for_status()
    return r.json().get("response", "")


def get_weather():
    try:
        r = requests.get("https://wttr.in/Seoul?format=3", timeout=20)
        r.raise_for_status()
        return r.text
    except Exception as e:
        return f"날씨 정보를 가져오지 못했습니다: {e}"


def select_aurum_route(question):
    q = question.lower()

    if any(k in q for k in ["날씨", "기온", "비", "눈", "weather"]):
        return "WEATHER_API", "날씨 질문으로 판단"

    if any(k in q for k in ["이미지", "사진", "ocr", "그림", "캡처"]):
        return MODEL_IMAGE, "이미지/OCR 관련 질문"

    if any(k in q for k in ["python", "코드", "에러", "오류", "git", "docker", "vscode", "터미널"]):
        return MODEL_FAST, "코드/시스템 질문"

    if any(k in q for k in ["요약", "문서", "보고서", "분석", "ppt", "엑셀", "환경영향평가"]):
        return MODEL_BIG, "문서/보고서 분석"

    return MODEL_LIGHT, "일반 질문"


def aurum_answer(question):
    if not question or not question.strip():
        return "선택 없음", "질문을 입력하세요."

    route, reason = select_aurum_route(question)

    if route == "WEATHER_API":
        return f"[AURUM]\n선택 도구: WEATHER_API\n선택 이유: {reason}", get_weather()

    prompt = f"""
너는 AURUM(아우룸) 통합 AI 시스템이다.
DGX 로컬 서버에서 OCR, Excel, CAD, KML, PDF, PPT 자동화를 지원한다.

선택 모델: {route}
선택 이유: {reason}

답변 원칙:
- 한국어
- 실무 중심
- 초보자도 따라할 수 있게 설명
- 명령어는 복붙 가능하게 제공
- 불확실하면 단정하지 않기

사용자 질문:
{question}
"""
    answer = call_ollama(route, prompt)
    return f"[AURUM]\n선택 모델: {route}\n선택 이유: {reason}", answer


def ask_llm_for_ocr(text):
    prompt = f"""
너는 회사 업무용 문서 분석 전문가다.
아래 OCR 결과는 여러 OCR 엔진 결과가 합쳐진 원문일 수 있다.

중요 지시:
- [TESSERACT_RESULT]와 [EASYOCR_RESULT]가 있으면 둘을 비교해서 최종 보정본을 먼저 만들어라.
- OCR 오타는 문맥상 가능한 범위에서 보정해라.
- 확실하지 않은 내용은 "추정"이라고 표시해라.
- 결과는 한국어 실무 보고서 형식으로 작성해라.

[최종 보정본]

[문서 유형]

[핵심 요약]
1.
2.
3.

[상세 분석]

[실행할 작업]
1.
2.
3.

[주의할 점]

OCR 원문:
{text}
"""
    return call_ollama(MODEL_BIG, prompt, timeout=600)


def create_report_ppt(title, ocr_text, ai_result, out_path):
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "AURUM 자동 리포트"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "OCR 원문"
    slide.placeholders[1].text = ocr_text[:1000]

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "AI 분석"
    slide.placeholders[1].text = ai_result[:1000]

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "다음 작업"
    slide.placeholders[1].text = "1. OCR 보정본 확인\n2. AI 분석 검토\n3. 필요 시 원본 개선\n4. 리포트 저장 및 공유"

    prs.save(out_path)


def run_ocr_ai_ppt(file):
    src = Path(file)
    path = BASE / "uploads" / src.name
    shutil.copy(src, path)

    subprocess.run(["python", str(BASE / "scripts/ocr_run.py"), str(path)], check=True)

    ocr_txt_path = BASE / "outputs" / f"{path.stem}_ocr.txt"
    ocr_img_path = BASE / "outputs" / f"{path.stem}_processed.png"

    ocr_text = ocr_txt_path.read_text(encoding="utf-8")
    ai_result = ask_llm_for_ocr(ocr_text)

    ai_file = BASE / "outputs" / f"{path.stem}_ai.txt"
    ai_file.write_text(ai_result, encoding="utf-8")

    ppt_file = BASE / "outputs" / f"{path.stem}_report.pptx"
    create_report_ppt("AURUM OCR 분석 리포트", ocr_text, ai_result, ppt_file)

    return str(ocr_txt_path), str(ocr_img_path), ai_result, str(ai_file), str(ppt_file)


def run_taxa_image_ocr(file):
    src = Path(file)
    path = BASE / "uploads" / src.name
    shutil.copy(src, path)

    subprocess.run(
        ["python", str(BASE / "scripts/taxa_table_ocr.py"), str(path)],
        check=True,
    )

    out = BASE / "outputs" / f"{path.stem}_taxa.xlsx"
    return str(out)


def safe_page_name(pages):
    p = pages.strip() if pages and pages.strip() else "all"
    return p.replace(",", "_").replace("-", "to").replace(" ", "")


def run_taxa_pdf_ocr(file, pages):
    src = Path(file)
    path = BASE / "uploads" / src.name
    shutil.copy(src, path)

    page_text = pages.strip() if pages and pages.strip() else ""

    subprocess.run(
        ["python", str(BASE / "scripts/taxa_pdf_ocr.py"), str(path), page_text],
        check=True,
    )

    out_name = f"{path.stem}_pages_{safe_page_name(page_text)}_taxa.xlsx"
    out = BASE / "outputs" / out_name

    return str(out)


def run_excel(file):
    src = Path(file)
    path = BASE / "uploads" / src.name
    shutil.copy(src, path)

    subprocess.run(["python", str(BASE / "scripts/excel_summary.py"), str(path)], check=True)

    out = BASE / "outputs" / f"{path.stem}_summary.xlsx"
    return str(out)


def auto_convert_cad_kml(file):
    src = Path(file)
    path = BASE / "uploads" / src.name
    shutil.copy(src, path)

    ext = path.suffix.lower()

    if ext == ".kml":
        mode = "kml_to_dxf"
        out = BASE / "outputs" / f"{path.stem}.dxf"
    elif ext == ".dxf":
        mode = "dxf_to_kml"
        out = BASE / "outputs" / f"{path.stem}.kml"
    else:
        raise ValueError("지원하지 않는 파일입니다. .kml 또는 .dxf 파일만 업로드하세요.")

    result = subprocess.run(
        ["python", str(BASE / "scripts/cad_kml_convert.py"), mode, str(path), str(out)],
        capture_output=True,
        text=True,
    )

    if result.returncode != 0:
        raise RuntimeError(result.stderr)

    log = f"변환 완료: {path.name} → {out.name}\n\n{result.stdout}"
    return str(out), log


def aurum_auto(file):
    src = Path(file)
    ext = src.suffix.lower()

    try:
        if ext in [".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"]:
            ocr_txt, ocr_img, ai_text, ai_file, ppt_file = run_ocr_ai_ppt(file)
            log = f"이미지 파일 감지: {src.name}\n처리: OCR → AI 분석 → PPT 생성 완료"
            return log, ocr_txt, ocr_img, ai_text, ai_file, ppt_file, None

        if ext in [".xlsx", ".xls"]:
            excel_out = run_excel(file)
            log = f"엑셀 파일 감지: {src.name}\n처리: Excel 분석 완료"
            return log, None, None, "", None, None, excel_out

        if ext in [".kml", ".dxf"]:
            converted, cad_log = auto_convert_cad_kml(file)
            log = f"CAD/KML 파일 감지: {src.name}\n{cad_log}"
            return log, None, None, "", None, None, converted

        if ext == ".pdf":
            log = f"PDF 파일 감지: {src.name}\nPDF는 페이지 지정이 필요하므로 'PDF 분류군 조사표 OCR' 탭에서 처리하세요."
            return log, None, None, "", None, None, None

        raise ValueError("지원하지 않는 파일입니다. 이미지, 엑셀, KML, DXF, PDF만 지원합니다.")

    except Exception as e:
        return f"오류 발생: {e}", None, None, "", None, None, None


def run_ppt():
    subprocess.run(["python", str(BASE / "scripts/ppt_create.py")], check=True)
    return str(BASE / "outputs" / "test_presentation.pptx")


def run_cad_test():
    subprocess.run(["python", str(BASE / "scripts/cad_create_dxf.py")], check=True)
    return str(BASE / "outputs" / "test_drawing.dxf")


with gr.Blocks(title="AURUM") as app:
    gr.Markdown("# AURUM / 아우룸")
    gr.Markdown("### DGX 통합 AI 작업 콘솔")

    with gr.Tab("AURUM AUTO"):
        gr.Markdown("## 파일 하나를 업로드하면 AURUM이 자동으로 처리합니다.")
        auto_file = gr.File(label="파일 업로드", type="filepath")
        auto_btn = gr.Button("AURUM AUTO 실행")

        auto_log = gr.Textbox(label="처리 로그", lines=6)
        auto_ocr_txt = gr.File(label="OCR 결과")
        auto_ocr_img = gr.Image(label="전처리 이미지")
        auto_ai_text = gr.Textbox(label="AI 분석 결과", lines=12)
        auto_ai_file = gr.File(label="AI 분석 파일")
        auto_ppt_file = gr.File(label="PPT 리포트")
        auto_result_file = gr.File(label="기타 결과 파일")

        auto_btn.click(
            aurum_auto,
            auto_file,
            [auto_log, auto_ocr_txt, auto_ocr_img, auto_ai_text, auto_ai_file, auto_ppt_file, auto_result_file],
        )

    with gr.Tab("AURUM 자동 질문"):
        question = gr.Textbox(label="질문 입력", lines=8)
        ask_btn = gr.Button("AURUM 실행")
        model_info = gr.Textbox(label="선택된 모델 / 도구", lines=4)
        answer = gr.Textbox(label="답변", lines=18)
        ask_btn.click(aurum_answer, question, [model_info, answer])

    with gr.Tab("Open WebUI"):
        gr.Markdown("## Open WebUI")
        gr.Markdown(f"[Open WebUI 새 창에서 열기]({OPEN_WEBUI_URL})")
        gr.HTML(f"""
        <iframe src="{OPEN_WEBUI_URL}" width="100%" height="800"
        style="border:1px solid #ccc; border-radius:8px;"></iframe>
        """)

    with gr.Tab("OCR → AI → PPT"):
        file = gr.File(label="이미지 업로드", type="filepath")
        btn = gr.Button("OCR + AI 분석 + PPT 생성")
        ocr_txt = gr.File(label="OCR 결과")
        ocr_img = gr.Image(label="전처리 이미지")
        ai_text = gr.Textbox(label="AI 분석", lines=12)
        ai_file = gr.File(label="AI 결과 파일")
        ppt_file = gr.File(label="PPT 리포트")
        btn.click(run_ocr_ai_ppt, file, [ocr_txt, ocr_img, ai_text, ai_file, ppt_file])

    with gr.Tab("분류군 조사표 OCR"):
        taxa_file = gr.File(label="조사표 이미지 업로드", type="filepath")
        taxa_btn = gr.Button("이미지 표 OCR → Excel 생성")
        taxa_out = gr.File(label="결과 Excel")
        taxa_btn.click(run_taxa_image_ocr, taxa_file, taxa_out)

    with gr.Tab("PDF 분류군 조사표 OCR"):
        pdf_file = gr.File(label="PDF 업로드", type="filepath")
        page_input = gr.Textbox(
            label="OCR할 페이지",
            value="1",
            placeholder="예: 1,3-5,8 / 전체 페이지는 빈칸"
        )
        pdf_btn = gr.Button("지정 페이지 OCR → Excel 생성")
        pdf_out = gr.File(label="결과 Excel")
        pdf_btn.click(run_taxa_pdf_ocr, [pdf_file, page_input], pdf_out)

    with gr.Tab("Excel"):
        excel_file = gr.File(label="엑셀 파일 업로드", type="filepath")
        excel_btn = gr.Button("엑셀 분석 실행")
        excel_out = gr.File(label="엑셀 분석 결과")
        excel_btn.click(run_excel, excel_file, excel_out)

    with gr.Tab("CAD / KML 자동 변환"):
        cad_kml_file = gr.File(label="KML 또는 DXF 파일 업로드", type="filepath")
        cad_kml_btn = gr.Button("자동 변환")
        cad_kml_out = gr.File(label="변환 결과 파일")
        cad_kml_log = gr.Textbox(label="처리 로그", lines=6)
        cad_kml_btn.click(auto_convert_cad_kml, cad_kml_file, [cad_kml_out, cad_kml_log])

    with gr.Tab("PPT 테스트"):
        ppt_btn = gr.Button("PPT 생성")
        ppt_out = gr.File(label="생성된 PPT")
        ppt_btn.click(run_ppt, None, ppt_out)

    with gr.Tab("CAD 테스트"):
        cad_btn = gr.Button("DXF 테스트 도면 생성")
        cad_out = gr.File(label="생성된 DXF")
        cad_btn.click(run_cad_test, None, cad_out)


app.launch(server_name="0.0.0.0", server_port=7861)
